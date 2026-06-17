"""Sections → PPTX exporter.

Input shape:
  {
    "title": "Deck Title",
    "subtitle": "optional",
    "sections": [
      {"heading": "Topic", "bullets": ["one", "two"], "notes": "optional"},
      ...
    ],
  }

Layouts used:
  - slide 0: title slide (layout 0)
  - slide N: title + content (layout 1) — heading + bulleted body

`build_pptx(title, sections, out_path, subtitle=None)` writes a .pptx
and returns the absolute path. Never raises on weird input — empty
bullet lists become a single "(no content)" placeholder so the deck
still renders.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt


def build_pptx(
    title: str,
    sections: list[dict[str, Any]],
    out_path: Path,
    *,
    subtitle: str | None = None,
) -> Path:
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    pres = Presentation()
    title_layout = pres.slide_layouts[0]
    content_layout = pres.slide_layouts[1]

    # Title slide.
    s = pres.slides.add_slide(title_layout)
    s.shapes.title.text = title or "Untitled deck"
    if subtitle and len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle

    for section in sections or []:
        if not isinstance(section, dict):
            continue
        heading = str(section.get("heading", "")).strip() or "Section"
        bullets = section.get("bullets", []) or []
        if isinstance(bullets, str):
            bullets = [bullets]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        if not bullets:
            bullets = ["(no content)"]
        notes = str(section.get("notes", "")).strip()

        slide = pres.slides.add_slide(content_layout)
        slide.shapes.title.text = heading[:120]

        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is None:
            # Create a fallback text box if the layout has no body
            # placeholder (some templates ship without one).
            body = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(5),
            )
        tf = body.text_frame
        tf.word_wrap = True
        tf.text = bullets[0][:300]
        for line in bullets[1:]:
            p = tf.add_paragraph()
            p.text = line[:300]
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes[:2000]

    pres.save(str(out_path))
    return out_path.resolve()


__all__ = ["build_pptx"]
